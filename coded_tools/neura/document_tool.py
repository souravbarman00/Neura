"""Neura CodedTool: create a downloadable PDF or PPTX from markdown-ish content.

The file is written under data/artifacts/ (served by the backend) and returned as a
markdown link, which the chat renders as a download card. Content is plain text with
simple structure:
  # / ## headings   → PPTX: start a new slide (its title);  PDF: a heading
  - / * lines       → bullet points
  blank line        → paragraph break
"""
from __future__ import annotations

import asyncio
import re
import time
import uuid
from pathlib import Path
from typing import Any, Dict, List, Tuple

from neuro_san.interfaces.coded_tool import CodedTool

ROOT = Path(__file__).resolve().parents[2]
ARTIFACTS = ROOT / "data" / "artifacts"


def _slug(s: str) -> str:
    s = re.sub(r"[^a-zA-Z0-9]+", "-", (s or "document")).strip("-").lower()
    return (s[:48] or "document")


def _parse(content: str) -> List[Tuple[str, List[str]]]:
    """Group content into (heading, lines[]) sections by # / ## headings."""
    sections: List[Tuple[str, List[str]]] = []
    cur_head = ""
    cur_lines: List[str] = []
    for raw in (content or "").splitlines():
        line = raw.rstrip()
        m = re.match(r"^\s{0,3}#{1,3}\s+(.*)$", line)
        if m:
            if cur_head or cur_lines:
                sections.append((cur_head, cur_lines))
            cur_head, cur_lines = m.group(1).strip(), []
        else:
            cur_lines.append(line)
    if cur_head or cur_lines:
        sections.append((cur_head, cur_lines))
    return sections or [("", [content or ""])]


def _bullet(line: str) -> Tuple[bool, str]:
    m = re.match(r"^\s*[-*]\s+(.*)$", line)
    return (True, m.group(1)) if m else (False, line.strip())


def _make_pptx(path: Path, title: str, content: str) -> None:
    from pptx import Presentation
    from pptx.util import Pt

    prs = Presentation()
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_slide.shapes.title.text = title or "Presentation"

    for head, lines in _parse(content):
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = head or title or "Slide"
        body = slide.placeholders[1].text_frame
        body.clear()
        first = True
        for ln in lines:
            if not ln.strip():
                continue
            is_b, text = _bullet(ln)
            p = body.paragraphs[0] if first else body.add_paragraph()
            first = False
            p.text = text
            p.level = 0 if is_b else 0
            p.font.size = Pt(18)
    prs.save(str(path))


def _make_pdf(path: Path, title: str, content: str) -> None:
    from fpdf import FPDF
    from fpdf.enums import XPos, YPos

    def safe(t: str) -> str:  # core fonts are latin-1; replace anything else
        return t.encode("latin-1", "replace").decode("latin-1")

    pdf = FPDF()
    pdf.set_auto_page_break(auto=True, margin=15)
    pdf.add_page()

    def cell(h: float, text: str) -> None:
        # new_x=LMARGIN resets the cursor to the left margin so the next full-width
        # multi_cell always has room (avoids fpdf2's "not enough space" error).
        pdf.multi_cell(0, h, safe(text), new_x=XPos.LMARGIN, new_y=YPos.NEXT)

    pdf.set_font("Helvetica", "B", 20)
    cell(10, title or "Document")
    pdf.ln(2)
    for head, lines in _parse(content):
        if head:
            pdf.set_font("Helvetica", "B", 15)
            cell(8, head)
            pdf.ln(1)
        pdf.set_font("Helvetica", "", 11)
        for ln in lines:
            if not ln.strip():
                pdf.ln(3)
                continue
            is_b, text = _bullet(ln)
            cell(6, ("  -  " + text) if is_b else text)
    pdf.output(str(path))


class MakeDocument(CodedTool):
    """Generate a downloadable PDF or PPTX from a title + markdown-ish content."""

    def invoke(self, args: Dict[str, Any], sly_data: Dict[str, Any]) -> Any:
        fmt = (args.get("format") or "").strip().lower()
        if fmt in ("ppt", "powerpoint", "slides"):
            fmt = "pptx"
        if fmt not in ("pdf", "pptx"):
            return "Set `format` to 'pdf' or 'pptx'."
        title = (args.get("title") or "").strip()
        content = args.get("content") or ""
        if not str(content).strip():
            return "Provide `content` (markdown-ish text) for the document."

        ARTIFACTS.mkdir(parents=True, exist_ok=True)
        name = f"{_slug(title)}-{int(time.time())}-{uuid.uuid4().hex[:4]}.{fmt}"
        out = ARTIFACTS / name
        try:
            if fmt == "pptx":
                _make_pptx(out, title, content)
            else:
                _make_pdf(out, title, content)
        except Exception as exc:  # noqa: BLE001
            return f"Could not create the {fmt.upper()}: {exc}"
        label = title or name
        return f"[📄 {label} ({fmt.upper()})](/artifacts/{name})"

    async def async_invoke(self, args: Dict[str, Any], sly_data: Dict[str, Any]) -> Any:
        return await asyncio.to_thread(self.invoke, args, sly_data)
